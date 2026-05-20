"""
Generate PowerPoint: Toxoplasma chromatin remodeler 2025 four-paper comparison
Author: Doctor Shang — research notes, May 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)   # slide background / title bar
MID_BLUE    = RGBColor(0x2E, 0x6D, 0xA8)   # accent / table header
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # table row tint
ORANGE      = RGBColor(0xE8, 0x7A, 0x22)   # highlight / conflict marker
GREEN       = RGBColor(0x27, 0x7A, 0x3B)   # positive / tick
RED         = RGBColor(0xC0, 0x39, 0x2B)   # negative / conflict
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xF8)
DARK_GREY   = RGBColor(0x2C, 0x3E, 0x50)
GOLD        = RGBColor(0xF3, 0x9C, 0x12)

PAPER_COLORS = {
    "Pachano": RGBColor(0x2E, 0x86, 0xAB),
    "Kashyap": RGBColor(0xA2, 0x3B, 0x72),
    "Zhu":     RGBColor(0xF1, 0x87, 0x01),
    "Hu":      RGBColor(0x2D, 0x93, 0x5F),
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_w:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, size=14, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=12, bold=False, italic=False,
             color=DARK_GREY, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None, bar_color=DARK_BLUE):
    """Dark top bar with title and optional subtitle."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=bar_color)
    add_textbox(slide, 0.35, 0.08, 12.5, 0.65, title,
                size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.35, 0.72, 12.5, 0.38, subtitle,
                    size=13, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # bottom rule
    add_rect(slide, 0, 7.3, 13.33, 0.1, fill_rgb=bar_color)
    # footer text
    add_textbox(slide, 0.35, 7.28, 12.5, 0.22,
                "Doctor Shang · Research Notes · May 2026",
                size=9, color=RGBColor(0x99,0xAA,0xBB), align=PP_ALIGN.LEFT)

def bullet_box(slide, l, t, w, h, items, size=12.5, color=DARK_GREY,
               bullet="▸", title=None, title_color=DARK_BLUE):
    """A text box with optional title + bullet list."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def simple_table(slide, l, t, col_widths, rows,
                 header_fill=MID_BLUE, alt_fill=LIGHT_BLUE,
                 font_size=10.5, header_font_size=11):
    """Draw a table manually as coloured rectangles + text boxes."""
    row_h = 0.38
    for ri, row in enumerate(rows):
        is_header = (ri == 0)
        fill = header_fill if is_header else (LIGHT_BLUE if ri % 2 == 0 else WHITE)
        txt_color = WHITE if is_header else DARK_GREY
        x = l
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, t + ri * row_h, cw, row_h,
                     fill_rgb=fill,
                     line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=0.5)
            add_textbox(slide, x + 0.05, t + ri * row_h + 0.03,
                        cw - 0.1, row_h - 0.05,
                        str(cell),
                        size=header_font_size if is_header else font_size,
                        bold=is_header,
                        color=txt_color,
                        align=PP_ALIGN.LEFT)
            x += cw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
# gradient-ish background via two rects
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(sl, 0, 4.8, 13.33, 2.7, fill_rgb=RGBColor(0x0F,0x23,0x3A))

# decorative bar
add_rect(sl, 0, 4.6, 13.33, 0.18, fill_rgb=ORANGE)

# paper colour strips at bottom
for i, (name, clr) in enumerate(PAPER_COLORS.items()):
    add_rect(sl, i * 3.33, 7.28, 3.33, 0.22, fill_rgb=clr)
    add_textbox(sl, i * 3.33 + 0.1, 7.27, 3.2, 0.25,
                name + " et al. 2025", size=9, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, 0.8, 0.7, 11.5, 1.1,
            "Toxoplasma gondii 染色质重塑器",
            size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 1.7, 11.5, 0.9,
            "2025 年四篇竞争论文横向比较",
            size=28, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 2.7, 11.5, 0.6,
            "Pachano · Kashyap · Zhu · Hu",
            size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(sl, 0.8, 4.95, 11.5, 0.45,
            "Nature Microbiology  |  Nature Communications  ×3",
            size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 5.45, 11.5, 0.45,
            "April – November 2025",
            size=13, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 6.2, 11.5, 0.45,
            "Doctor Shang · Research Notes · May 2026",
            size=11, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  — Background & Landscape
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "背景：Toxoplasma 染色质调控全景",
             subtitle="为什么 2025 年突然出现 4 篇相关论文？")

# Left panel — key facts
add_rect(sl, 0.3, 1.3, 5.9, 5.7, fill_rgb=WHITE,
         line_rgb=MID_BLUE, line_w=1.5)
add_rect(sl, 0.3, 1.3, 5.9, 0.45, fill_rgb=MID_BLUE)
add_textbox(sl, 0.35, 1.32, 5.8, 0.4,
            "T. gondii 基因组关键事实", size=12, bold=True,
            color=WHITE)

facts = [
    "75% 基因含内含子（vs 酵母 5%）— intron-rich genome",
    "复杂生活史：速殖子 → 缓殖子 → 裂殖子 → 配子体 → 卵囊",
    "每个生活史阶段有独立的基因表达程序",
    "67 个 ApiAP2 转录因子负责序列特异性结合",
    "MORC + HDAC3 是中心沉默枢纽，与 ≥12 个 ApiAP2 互作",
]
bullet_box(sl, 0.4, 1.85, 5.7, 4.9, facts, size=12, color=DARK_GREY,
           bullet="◆")

# Right panel — 4 remodeler families
add_rect(sl, 6.5, 1.3, 6.5, 5.7, fill_rgb=WHITE,
         line_rgb=DARK_BLUE, line_w=1.5)
add_rect(sl, 6.5, 1.3, 6.5, 0.45, fill_rgb=DARK_BLUE)
add_textbox(sl, 6.55, 1.32, 6.4, 0.4,
            "四大类 ATP 依赖染色质重塑器", size=12, bold=True,
            color=WHITE)

families = [
    ("ISWI", "TgSNF2h (TGME49_321440)\nTgSNF2L (TGME49_273870)", PAPER_COLORS["Pachano"]),
    ("SWI/SNF", "TgSNF2a (TGME49_278440)\nTgSNF2b (TGME49_320300)", PAPER_COLORS["Hu"]),
    ("INO80", "TgSRCAP + TGME49_229460\nTGME49_226440", MID_BLUE),
    ("CHD", "TgCHD1 (TGME49_258240)", DARK_GREY),
]
for i, (fam, members, clr) in enumerate(families):
    y = 1.9 + i * 1.2
    add_rect(sl, 6.6, y, 1.3, 0.95, fill_rgb=clr)
    add_textbox(sl, 6.62, y + 0.22, 1.26, 0.5, fam,
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 7.95, y, 4.9, 0.95, fill_rgb=LIGHT_BLUE,
             line_rgb=clr, line_w=1)
    add_textbox(sl, 8.0, y + 0.1, 4.8, 0.8, members,
                size=11, color=DARK_GREY)

# Timeline bar at bottom
add_rect(sl, 0.3, 6.8, 12.73, 0.35, fill_rgb=DARK_BLUE)
add_textbox(sl, 0.35, 6.82, 12.6, 0.28,
            "2025  Apr 11 Pachano  ·  Apr 22 Kashyap  ·  Jul 01 Zhu  ·  Nov 05 Hu",
            size=11, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  — Four Papers Overview Table
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "四篇论文基本信息",
             subtitle="同一生物问题的四个切入角度，7 个月内密集发表")

headers = ["项目", "Pachano et al.", "Kashyap et al.", "Zhu et al.", "Hu et al."]
rows = [
    headers,
    ["期刊", "Nat Microbiol", "Nat Commun", "Nat Commun", "Nat Commun"],
    ["卷期", "10:1156-1170", "16:3769", "16:5757", "16:9777"],
    ["发表日期", "2025-04-11", "2025-04-22", "2025-07-01", "2025-11-05"],
    ["通讯作者", "M.-A. Hakimi", "A.S. Deshmukh", "B. Shen (沈彬)", "X. Song (宋兴军)"],
    ["实验室", "Grenoble, France", "Hyderabad, India", "武汉，华中农大", "南宁，广西大学"],
    ["研究对象", "TgSNF2h (+SNF2L对比)", "TgCdc5 剪接因子", "TgSNF2L 复合物", "TgSNF2a + TgSNF2b"],
    ["蛋白家族", "ISWI", "Spliceosome NTC", "ISWI", "SWI/SNF"],
    ["一句话核心", "Insulator 阻止\nMORC 蔓延", "剪接缺陷→\n流产型缓殖子", "SNF2L 抑制\n有性期基因", "SWI/SNF 激活\n速殖子基因"],
]
col_w = [2.0, 2.6, 2.6, 2.6, 2.6]
simple_table(sl, 0.3, 1.25, col_w, rows, font_size=10, header_font_size=11)

# colour labels for each paper
for i, (name, clr) in enumerate(PAPER_COLORS.items()):
    add_rect(sl, 2.3 + i * 2.6, 1.25, 2.6, 0.12, fill_rgb=clr)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  — Pachano et al.
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "Pachano et al. (2025) — TgSNF2h / ISWI 家族",
             subtitle='Nat Microbiol 10:1156 | "TgSNF2h is the insulator that keeps MORC from silencing active genes"',
             bar_color=PAPER_COLORS["Pachano"])

# core finding box
add_rect(sl, 0.3, 1.3, 12.73, 0.75, fill_rgb=PAPER_COLORS["Pachano"])
add_textbox(sl, 0.4, 1.35, 12.5, 0.65,
            "核心发现：TgSNF2h + TgRFTS + AP2VIII-2 形成 ISWI 复合物，在高转录基因上充当 Insulator，阻止 MORC 沉默机器向活跃区域蔓延",
            size=13, bold=True, color=WHITE)

# Left: innovations
add_rect(sl, 0.3, 2.15, 5.9, 4.75, fill_rgb=WHITE,
         line_rgb=PAPER_COLORS["Pachano"], line_w=1.5)
bullet_box(sl, 0.4, 2.2, 5.7, 4.6,
           ["首次系统鉴定 Toxoplasma 两个 ISWI ATPase",
            "发现 TgRFTS 作为 scaffold（SpRAF2 同源，但功能重定向）",
            "首提 Toxoplasma 染色质 'Insulator' 概念",
            "Epistatic control over MORC 蔓延",
            "TgRFTS KD 完美表型复现（phenocopy）TgSNF2h KD",
            "AlphaFold v2 + Foldseek 远缘结构同源分析",
            "Nanopore direct RNA-seq 发现 read-through"],
           title="核心创新点", size=11.5,
           title_color=PAPER_COLORS["Pachano"])

# Right: key data
add_rect(sl, 6.5, 2.15, 6.5, 4.75, fill_rgb=WHITE,
         line_rgb=PAPER_COLORS["Pachano"], line_w=1.5)
kd_rows = [
    ["实验", "结果"],
    ["TgSNF2h IP-MS", "富集 TgRFTS + AP2VIII-2"],
    ["TgSNF2h KD ATAC-seq", "Chromatin accessibility 显著降低"],
    ["TgSNF2L KD ATAC-seq", "⚠ 无显著变化（被 Zhu 挑战）"],
    ["TgSNF2h KD MORC ChIP", "MORC 蔓延到 active gene 区域"],
    ["双向启动子分析", "KD 后两基因反相关—支持 insulator 模型"],
    ["TgRFTS vs TgSNF2h KD", "转录组 & ATAC-seq 高度一致（phenocopy）"],
]
simple_table(sl, 6.55, 2.2, [2.5, 3.85], kd_rows,
             header_fill=PAPER_COLORS["Pachano"], font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  — Kashyap et al.
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "Kashyap et al. (2025) — TgCdc5 / 剪接因子",
             subtitle='Nat Commun 16:3769 | "TgCdc5 splicing maintains tachyzoite transcriptome; loss causes abortive bradyzoite"',
             bar_color=PAPER_COLORS["Kashyap"])

add_rect(sl, 0.3, 1.3, 12.73, 0.75, fill_rgb=PAPER_COLORS["Kashyap"])
add_textbox(sl, 0.4, 1.35, 12.5, 0.65,
            "核心发现：TgCdc5 是剪接体 NTC 核心；缺失导致广泛内含子滞留 → 蛋白错误折叠 → 自发激活流产型缓殖子分化 → 虫体死亡；小鼠 KD 可作减毒疫苗候选",
            size=13, bold=True, color=WHITE)

add_rect(sl, 0.3, 2.15, 5.9, 4.75, fill_rgb=WHITE,
         line_rgb=PAPER_COLORS["Kashyap"], line_w=1.5)
bullet_box(sl, 0.4, 2.2, 5.7, 4.6,
           ["首次在 Toxoplasma 中表征剪接体 NTC 复合物",
            "确立剪接抑制 → 缓殖子诱导因果链",
            "酵母功能互补实验证明进化保守性",
            "小鼠感染：KD 株完全减毒，激发保护性免疫",
            "妊娠小鼠部分保护数据",
            "Splicing minigene reporter 系统",
            "MST 测 RNA-蛋白亲和力（U2 KD=440 nM）"],
           title="核心创新点", size=11.5,
           title_color=PAPER_COLORS["Kashyap"])

kd_rows = [
    ["实验", "结果"],
    ["TgCdc5 IP-MS", "146 互作蛋白，52 剪接体成员，NTC 8 核心"],
    ["RNA-seq 剪接分析", "20,987 事件受影响，82% 内含子滞留"],
    ["Splicing reporter", "KD 8h 后效率显著下降"],
    ["TUNEL / Annexin-PI", "~80% DNA 损伤；~32% 凋亡样"],
    ["Proteostat 染色", "~22% 蛋白聚集体"],
    ["小鼠感染", "KD 株减毒；激发保护性免疫"],
    ["酵母互补", "TgCdc5 拯救 Δcef1，TgPrp19 拯救 Δprp19"],
]
simple_table(sl, 6.55, 2.2, [2.5, 3.85], kd_rows,
             header_fill=PAPER_COLORS["Kashyap"], font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  — Zhu et al.
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "Zhu et al. (2025) — TgSNF2L / ISWI 家族",
             subtitle='Nat Commun 16:5757 | "TgSNF2L complex represses sexual genes; loss triggers in vitro out-budding"',
             bar_color=PAPER_COLORS["Zhu"])

add_rect(sl, 0.3, 1.3, 12.73, 0.75, fill_rgb=PAPER_COLORS["Zhu"])
add_textbox(sl, 0.4, 1.35, 12.5, 0.65,
            "核心发现：TgSNF2L + SLIF1 + SLIF2 + AP2X-4 形成 4 亚基复合物（类 yeast ISW2），抑制肠上皮期/有性期基因；KD 导致罕见体外 out-budding 分裂模式",
            size=13, bold=True, color=WHITE)

add_rect(sl, 0.3, 2.15, 5.9, 4.75, fill_rgb=WHITE,
         line_rgb=PAPER_COLORS["Zhu"], line_w=1.5)
bullet_box(sl, 0.4, 2.2, 5.7, 4.6,
           ["首次发现体外 out-budding 分裂模式（猫肠道特征）",
            "鉴定 SLIF1/SLIF2：Coccidia 特有的 accessory 蛋白",
            "系统筛选 16 个 putative ATPase",
            "直接挑战 Pachano：TgSNF2L KD 显著影响染色质可及性",
            "mAID 6h 完全降解（vs Pachano 48h 仍可见）",
            "TMT 蛋白组学 + 转录组学双重验证",
            "SFINX 程序精确分析 IP-MS 互作"],
           title="核心创新点", size=11.5,
           title_color=PAPER_COLORS["Zhu"])

kd_rows = [
    ["实验", "结果"],
    ["6 个候选 mAID 表型筛", "TgSNF2L KD 最严重（75% 异常分裂）"],
    ["WB 降解动力学", "IAA 6h 即不可检测（关键对照）"],
    ["分裂模式量化", "36h: 11% 裂殖子，38% out-budding"],
    ["RNA-seq 12h KD", "301 上调 / 193 下调；81% 上调 = EES/慢性期"],
    ["TMT 蛋白组", "GRA11B, AAH2, PAN, MIC17B 蛋白上调"],
    ["CUT&Tag", "91.7% peaks 在启动子；两 strain 80% 重叠"],
    ["4 亚基复合物", "SNF2L+SLIF1+SLIF2+AP2X-4 reciprocal 验证"],
]
simple_table(sl, 6.55, 2.2, [2.5, 3.85], kd_rows,
             header_fill=PAPER_COLORS["Zhu"], font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  — Hu et al.
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "Hu et al. (2025) — TgSNF2a + TgSNF2b / SWI/SNF 家族",
             subtitle='Nat Commun 16:9777 | "Two SWI/SNF ATPases cooperatively activate tachyzoite genes via direct MORC interaction"',
             bar_color=PAPER_COLORS["Hu"])

add_rect(sl, 0.3, 1.3, 12.73, 0.75, fill_rgb=PAPER_COLORS["Hu"])
add_textbox(sl, 0.4, 1.35, 12.5, 0.65,
            "核心发现：TgSNF2a 和 TgSNF2b 协同激活速殖子特异性细胞器基因；KD 导致不完全内复殖 + 淀粉蓄积；TgSNF2b 与 MORC/HDAC3 直接互作",
            size=13, bold=True, color=WHITE)

add_rect(sl, 0.3, 2.15, 5.9, 4.75, fill_rgb=WHITE,
         line_rgb=PAPER_COLORS["Hu"], line_w=1.5)
bullet_box(sl, 0.4, 2.2, 5.7, 4.6,
           ["首次完整描述 Toxoplasma SWI/SNF 复合物",
            "TgSNF2b 与 MORC/HDAC3 直接 Co-IP 互作",
            "TgSNF2a × TgSNF2b 直接互作（Fig. 7d）",
            "淀粉（amylopectin）蓄积表型（PAS 染色）—独特发现",
            "双 KD + ATAC-seq 证明协同染色质可及性降低",
            "TgMAPK-L1 机制解释不完全内复殖",
            "Astral DIA 高灵敏度定量蛋白组",
            "严格阶段标志物组合 (SAG1+GRA82+IMC7+DBL) 否定真分化"],
           title="核心创新点", size=11.5,
           title_color=PAPER_COLORS["Hu"])

kd_rows = [
    ["实验", "结果"],
    ["RNA-seq (Pearson)", "TgSNF2a vs TgSNF2b: r = 0.95，共享 2207 DEG"],
    ["多核率 (TgSNF2b)", "11% → 72%（第 2→7 天）"],
    ["TgSNF2a×b Co-IP", "直接互作（Fig.7d）"],
    ["TgSNF2b × MORC", "直接 Co-IP（Fig.7e）"],
    ["TgSNF2b × HDAC3", "直接 Co-IP（Fig.7f）"],
    ["ATAC-seq (dKD)", "455 启动子可及性降低，98% 下调"],
    ["MORC 结合重叠", "86% SWI/SNF 启动子与 MORC 重叠（P=9e-139）"],
]
simple_table(sl, 6.55, 2.2, [2.5, 3.85], kd_rows,
             header_fill=PAPER_COLORS["Hu"], font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  — Methodology Comparison
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "方法学全景对比",
             subtitle="四篇论文技术路线的优势与互补")

tech_rows = [
    ["技术 / 方法", "Pachano", "Kashyap", "Zhu", "Hu"],
    ["mAID 条件敲减", "✓ 慢(48h)", "✓ <1h", "✓ 6h", "✓ 12h"],
    ["IP-MS 互作组", "✓ LC-MS/MS", "✓", "✓", "✓ Astral DIA"],
    ["ChIP-seq", "✓", "—", "—", "—"],
    ["CUT&Tag", "—", "—", "✓", "✓"],
    ["ATAC-seq", "✓", "—", "✓", "✓ (dKD)"],
    ["Nanopore DRS", "✓ 独有", "—", "—", "—"],
    ["Illumina RNA-seq", "✓", "✓", "✓", "✓"],
    ["TMT/定量蛋白组", "—", "—", "✓", "✓"],
    ["酵母功能互补", "—", "✓ 独有", "—", "—"],
    ["MST RNA-蛋白亲和", "—", "✓ 独有", "—", "—"],
    ["TEM 超微结构", "—", "—", "✓", "✓"],
    ["小鼠体内感染", "—", "✓ 独有", "—", "—"],
    ["PAS 淀粉染色", "—", "—", "—", "✓ 独有"],
    ["AlphaFold 结构预测", "✓ 独有", "—", "—", "—"],
    ["阶段标志物 panel", "基础", "✓", "✓", "✓ 最全"],
]
col_w = [3.2, 2.3, 2.3, 2.3, 2.3]
simple_table(sl, 0.3, 1.25, col_w, tech_rows,
             header_fill=DARK_BLUE, font_size=10, header_font_size=11)

# score bars at bottom
add_textbox(sl, 0.35, 6.78, 5.0, 0.35,
            "方法全面性（综合评分）:", size=11, bold=True, color=DARK_GREY)
scores = [("Pachano", 8.5, PAPER_COLORS["Pachano"]),
          ("Kashyap", 8.0, PAPER_COLORS["Kashyap"]),
          ("Zhu",     7.0, PAPER_COLORS["Zhu"]),
          ("Hu",      8.0, PAPER_COLORS["Hu"])]
for i, (name, score, clr) in enumerate(scores):
    x = 5.0 + i * 2.0
    add_rect(sl, x, 6.75, score * 0.18, 0.38, fill_rgb=clr)
    add_textbox(sl, x + score * 0.18 + 0.05, 6.77, 0.8, 0.3,
                f"{name} {score}", size=10, color=clr, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9  — Core Findings Comparison Matrix
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "核心结论对比矩阵",
             subtitle="共识、分歧与各论文独有贡献")

# consensus box
add_rect(sl, 0.3, 1.3, 6.0, 2.95, fill_rgb=WHITE,
         line_rgb=GREEN, line_w=2)
add_rect(sl, 0.3, 1.3, 6.0, 0.4, fill_rgb=GREEN)
add_textbox(sl, 0.35, 1.32, 5.9, 0.36,
            "四篇共识（Consensus）", size=12, bold=True, color=WHITE)
bullet_box(sl, 0.4, 1.75, 5.8, 2.4,
           ["染色质重塑器 KD → 阶段特异性基因表达改变",
            "均出现异常分裂模式（内复殖样或裂殖子样）",
            "有性期 / EES 期基因的去抑制（de-repression）",
            "MORC-HDAC3 是沉默枢纽",
            "ApiAP2 转录因子在染色质靶向中关键",
            "Apicomplexa 特异性 scaffold/accessory 蛋白是新创新"],
           size=11, bullet="✔", color=DARK_GREY)

# conflict box
add_rect(sl, 6.6, 1.3, 6.4, 2.95, fill_rgb=WHITE,
         line_rgb=RED, line_w=2)
add_rect(sl, 6.6, 1.3, 6.4, 0.4, fill_rgb=RED)
add_textbox(sl, 6.65, 1.32, 6.3, 0.36,
            "⚠ 关键冲突", size=12, bold=True, color=WHITE)
conflict_rows = [
    ["议题", "Pachano", "Zhu"],
    ["TgSNF2L 影响\nATAC-seq？", "无影响", "显著降低"],
    ["TgSNF2L KD\n表型严重性", "轻微", "75% 异常"],
    ["根本原因", "—", "mAID 降解\n不完全"],
]
simple_table(sl, 6.65, 1.75, [2.5, 1.8, 1.8], conflict_rows,
             header_fill=RED, font_size=10)

# stage conversion comparison
add_rect(sl, 0.3, 4.4, 12.73, 2.75, fill_rgb=WHITE,
         line_rgb=ORANGE, line_w=2)
add_rect(sl, 0.3, 4.4, 12.73, 0.4, fill_rgb=ORANGE)
add_textbox(sl, 0.35, 4.42, 12.6, 0.36,
            "阶段转化（Stage Conversion）议题 — 各篇结论不同", size=12, bold=True, color=WHITE)
sc_rows = [
    ["论文", "观察到的表型", "是否真正阶段转化？", "机制解释"],
    ["Kashyap", "BAG1+, CST1+, DBL+ (缓殖子标志)", "流产型（abortive），虫体死亡", "剪接缺陷 → 应激 → 缓殖子信号"],
    ["Zhu",     "PF16+, 多线粒体（雄配子特征）", "部分形态变化，非完整分化", "有性期基因释放 → 形态改变"],
    ["Hu",      "裂殖子基因上调，SAG1+ 贯穿始终", "明确否定（Not conversion）", "去抑制 ≠ 阶段转化，速殖子身份保留"],
]
simple_table(sl, 0.35, 4.85, [1.8, 3.5, 3.5, 3.7], sc_rows,
             header_fill=ORANGE, font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10  — Key Conflicts Detail
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "关键冲突点详解",
             subtitle="Pachano vs Zhu (TgSNF2L) · Hu vs Schwarz (SWI/SNF architecture)")

# Conflict 1
add_rect(sl, 0.3, 1.25, 12.73, 0.4, fill_rgb=RED)
add_textbox(sl, 0.35, 1.27, 12.6, 0.36,
            "冲突 1：TgSNF2L KD 表型 — Pachano (2025-04) vs Zhu (2025-07)", size=12, bold=True, color=WHITE)

c1_rows = [
    ["指标", "Pachano et al.", "Zhu et al."],
    ["TgSNF2L KD 表型", "轻微", "75% 异常分裂（最严重）"],
    ["Chromatin accessibility", "无显著变化", "显著降低（54% target）"],
    ["mAID 降解验证", "48h WB 仍可检测到", "6h WB 即不可检测"],
    ["复合物描述", "未深挖", "SNF2L+SLIF1+SLIF2+AP2X-4"],
    ["结论", "TgSNF2h vs TgSNF2L 功能各异", "SNF2L 是重要的有性期基因抑制子"],
]
simple_table(sl, 0.35, 1.7, [3.5, 4.3, 4.3], c1_rows,
             header_fill=RED, font_size=10.5)

add_rect(sl, 0.35, 4.1, 12.63, 0.55, fill_rgb=RGBColor(0xFF,0xEE,0xEE),
         line_rgb=RED, line_w=1)
add_textbox(sl, 0.45, 4.13, 12.5, 0.48,
            'Zhu 原文："Such discrepancy is likely due to the incomplete depletion of SNF2L in the Pachano et al. study."  →  mAID 降解必须 WB 多时间点严格验证',
            size=11, italic=True, color=RED)

# Conflict 2
add_rect(sl, 0.3, 4.8, 12.73, 0.4, fill_rgb=PAPER_COLORS["Hu"])
add_textbox(sl, 0.35, 4.82, 12.6, 0.36,
            "冲突 2：SWI/SNF 复合物构型 — Hu (2025-11) vs Schwarz (2026-05)", size=12, bold=True, color=WHITE)

c2_rows = [
    ["模型", "Hu et al.", "Schwarz et al."],
    ["TgSNF2a × TgSNF2b", "直接 Co-IP 互作（Fig.7d）", "稳定性检验显示两个独立复合物"],
    ["解读", "协同双 ATPase 复合物", "两个互斥复合物"],
    ["证据基础", "直接生化（Co-IP）", "稳定性推断 + ALFA-CUT&RUN"],
    ["调和方式", "Dynamic exchange model — 共享 scaffold 但动态交换 ATPase", "需要 SEC-MALS/BN-PAGE/cryo-EM 最终裁决"],
]
simple_table(sl, 0.35, 5.25, [3.0, 4.5, 4.5], c2_rows,
             header_fill=PAPER_COLORS["Hu"], font_size=10.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11  — Unified Framework
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "统一框架：四篇 paper 收敛的核心模型",
             subtitle="三层调控模型 + 所有通路汇聚于速殖子身份维持")

# Three layers diagram
layer_data = [
    (1.3, "LAYER 1  ATP 依赖染色质重塑器", WHITE, DARK_BLUE,
     [("ISWI", "TgSNF2h+TgRFTS+AP2VIII-2\nInsulator，防 MORC 蔓延", PAPER_COLORS["Pachano"]),
      ("ISWI", "TgSNF2L+SLIF1+SLIF2+AP2X-4\n抑制有性期/EES 基因", PAPER_COLORS["Zhu"]),
      ("SWI/SNF", "TgSNF2a+TgSNF2b 复合物\n直接激活速殖子基因", PAPER_COLORS["Hu"])]),
    (3.3, "LAYER 2  沉默枢纽 (Silencing Hub)", WHITE, RED,
     [("MORC", "MORC + HDAC3\n+ AP2XII-1/XI-2/XII-5", RED),
      ("功能", "抑制阶段特异性基因\n在速殖子期 OFF", DARK_GREY),
      ("输入", "接收来自染色质重塑器\n的可及性信号", MID_BLUE)]),
    (5.1, "LAYER 3  转录组完整性", WHITE, PAPER_COLORS["Kashyap"],
     [("Splicing", "TgCdc5/NTC 剪接体\n确保正确 mRNA 产生", PAPER_COLORS["Kashyap"]),
      ("影响", "内含子滞留 → 错误蛋白\n→ 应激 → 缓殖子信号", ORANGE),
      ("结果", "功能蛋白翻译\n维持速殖子表型", GREEN)]),
]

for (y, title, txt_clr, bar_clr, items) in layer_data:
    add_rect(sl, 0.3, y, 12.73, 1.85, fill_rgb=WHITE,
             line_rgb=bar_clr, line_w=1.5)
    add_rect(sl, 0.3, y, 12.73, 0.35, fill_rgb=bar_clr)
    add_textbox(sl, 0.35, y + 0.03, 12.6, 0.3, title,
                size=12, bold=True, color=WHITE)
    for i, (sub, desc, clr) in enumerate(items):
        x = 0.5 + i * 4.1
        add_rect(sl, x, y + 0.42, 1.1, 1.3, fill_rgb=clr)
        add_textbox(sl, x + 0.05, y + 0.65, 1.0, 0.6,
                    sub, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, x + 1.15, y + 0.42, 2.8, 1.3,
                 fill_rgb=LIGHT_BLUE, line_rgb=clr, line_w=1)
        add_textbox(sl, x + 1.2, y + 0.5, 2.7, 1.1,
                    desc, size=10.5, color=DARK_GREY)

# outcome box
add_rect(sl, 0.3, 7.1, 12.73, 0.28, fill_rgb=DARK_BLUE)
add_textbox(sl, 0.35, 7.12, 12.6, 0.24,
            "OUTCOME ▸  速殖子身份维持 · 正确分裂方式 · 侵袭性基因 ON · 有性/缓殖子基因 OFF",
            size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12  — Open Questions & Future Directions
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
slide_header(sl, "开放问题与未来方向",
             subtitle="四篇 paper 留下的关键问题 + 方法学教训")

# Left: open questions table
add_rect(sl, 0.3, 1.3, 7.9, 5.7, fill_rgb=WHITE,
         line_rgb=DARK_BLUE, line_w=1.5)
add_rect(sl, 0.3, 1.3, 7.9, 0.4, fill_rgb=DARK_BLUE)
add_textbox(sl, 0.35, 1.32, 7.8, 0.36,
            "Open Questions", size=12, bold=True, color=WHITE)

oq_rows = [
    ["问题", "来源", "解决方案"],
    ["TgSNF2a vs SNF2b 是否\n形成两个独立复合物？", "Hu vs Schwarz", "SEC-MALS\nBN-PAGE\ncryo-EM"],
    ["TgRFTS RFTS domain\n在 Toxoplasma 的新功能？", "Pachano", "结构生物学\n+ domain swap"],
    ["SLIF1/SLIF2 在\nPlasmodium 的同源物？", "Zhu", "Plasmodium\nIP-MS 验证"],
    ["MORC 蔓延后的物理形式？\n是否形成 condensate？", "Pachano", "活细胞成像\n+ 生化表征"],
    ["不同重塑器 KD 释放\n不同基因子集的原因？", "全部", "全 KD panel +\n比较转录组"],
    ["细胞周期中重塑器动态？", "Schwarz 部分回答", "同步化 +\n时间分辨 CUT&RUN"],
]
simple_table(sl, 0.35, 1.75, [3.5, 2.0, 2.25], oq_rows,
             header_fill=MID_BLUE, font_size=10)

# Right: lessons
add_rect(sl, 8.5, 1.3, 4.5, 5.7, fill_rgb=WHITE,
         line_rgb=ORANGE, line_w=1.5)
add_rect(sl, 8.5, 1.3, 4.5, 0.4, fill_rgb=ORANGE)
add_textbox(sl, 8.55, 1.32, 4.4, 0.36,
            "方法学教训", size=12, bold=True, color=WHITE)

lessons = [
    "mAID 必须 WB 多时间点\n验证降解（Pachano vs Zhu）",
    "CUT&Tag 分辨率\n不及 CUT&RUN",
    "OE 系数据\n需谨慎解读（Hu）",
    "阶段标志物 panel\n不可或缺（Hu 示范）",
    "双 KD > 单 KD\n（功能冗余蛋白）",
    "De-repression ≠\nStage conversion",
]
for i, lesson in enumerate(lessons):
    y = 1.8 + i * 0.82
    add_rect(sl, 8.6, y, 0.35, 0.6, fill_rgb=ORANGE)
    add_textbox(sl, 8.62, y + 0.1, 0.3, 0.4,
                str(i+1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 9.0, y, 3.85, 0.6, fill_rgb=LIGHT_BLUE)
    add_textbox(sl, 9.05, y + 0.08, 3.75, 0.5,
                lesson, size=10.5, color=DARK_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13  — Summary / Take-home
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(sl, 0, 4.5, 13.33, 3.0, fill_rgb=RGBColor(0x0F,0x23,0x3A))
add_rect(sl, 0, 4.38, 13.33, 0.18, fill_rgb=ORANGE)

add_textbox(sl, 1.0, 0.4, 11.0, 0.8,
            "Take-Home Messages", size=28, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

messages = [
    ("MORC 是主协调者", "四篇 paper 中不同染色质重塑器均通过不同路径与 MORC-HDAC3 枢纽关联", PAPER_COLORS["Hu"]),
    ("De-repression ≠ 阶段转化", "Hu 严格区分了基因去抑制与真正阶段身份转换，是概念层面重要贡献", PAPER_COLORS["Zhu"]),
    ("应激 vs 染色质两条路径", "Kashyap（剪接）和 Hu（染色质）通过不同机制都可触发缓殖子信号", PAPER_COLORS["Kashyap"]),
    ("Apicomplexan-specific scaffolds", "TgRFTS, SLIF1/2 代表寄生虫特异的创新；是潜在的抗寄生虫靶点", PAPER_COLORS["Pachano"]),
]
for i, (title, desc, clr) in enumerate(messages):
    x = 0.4 + (i % 2) * 6.3
    y = 1.35 + (i // 2) * 1.55
    add_rect(sl, x, y, 6.0, 1.4, fill_rgb=RGBColor(0x1E,0x47,0x70),
             line_rgb=clr, line_w=2)
    add_rect(sl, x, y, 6.0, 0.35, fill_rgb=clr)
    add_textbox(sl, x + 0.1, y + 0.04, 5.8, 0.28,
                title, size=12, bold=True, color=WHITE)
    add_textbox(sl, x + 0.1, y + 0.42, 5.8, 0.85,
                desc, size=11, color=LIGHT_BLUE)

# four paper summary strips
for i, (name, clr) in enumerate(PAPER_COLORS.items()):
    x = 0.4 + i * 3.1
    add_rect(sl, x, 4.55, 3.0, 2.7, fill_rgb=RGBColor(0x1A,0x30,0x4A),
             line_rgb=clr, line_w=1.5)
    add_rect(sl, x, 4.55, 3.0, 0.35, fill_rgb=clr)
    labels = {
        "Pachano": ["Nat Microbiol", "TgSNF2h ISWI", "Insulator/MORC"],
        "Kashyap": ["Nat Commun", "TgCdc5 NTC", "Splicing→Bradyz."],
        "Zhu":     ["Nat Commun", "TgSNF2L ISWI", "Out-budding"],
        "Hu":      ["Nat Commun", "TgSNF2a/b SWI/SNF", "Direct MORC PPI"],
    }
    add_textbox(sl, x + 0.1, 4.57, 2.8, 0.28,
                f"{name} et al.", size=11, bold=True, color=WHITE)
    for j, line in enumerate(labels[name]):
        add_textbox(sl, x + 0.1, 5.0 + j * 0.42, 2.8, 0.38,
                    line, size=10.5, color=LIGHT_BLUE)

add_textbox(sl, 0.3, 7.28, 12.7, 0.22,
            "Doctor Shang · Research Notes · May 2026",
            size=9, color=RGBColor(0x66,0x88,0xAA), align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/inversion/Toxoplasma_chromatin_2025_comparison.pptx"
prs.save(out)
print(f"Saved → {out}")
